"""
FEAAM Priority Sender
=====================
1. Upload a leads CSV / Excel
2. Leads are scored and ranked (Priority 1 = 250-5000 employees, Priority 2 = 5000+)
3. Each lead gets a sector-matched deck auto-assigned as attachment
4. Review + edit the generated email in the browser
5. Click Send  →  email goes via Ionos SMTP with the deck attached
   (Case-study PDFs can be toggled on per email as secondary attachment)

All emails require MANUAL confirmation before sending — nothing goes out automatically.
"""

import io, os, re, json, uuid, smtplib, tempfile, threading
from email.message import EmailMessage
from datetime import datetime
from flask import Flask, render_template, jsonify, request, send_file
import pandas as pd

# ── Anthropic client (lazy — initialised after settings are loaded) ───────────
import anthropic as _anthropic

def _get_ai_client():
    """Return an Anthropic client using key from settings (or env as fallback)."""
    s   = load_settings()
    key = s.get("anthropic_api_key") or os.getenv("ANTHROPIC_API_KEY", "")
    if not key:
        return None, "No Anthropic API key set. Add it in Settings → Anthropic API Key."
    try:
        return _anthropic.Anthropic(api_key=key), None
    except Exception as e:
        return None, str(e)

# In-memory company size cache  {company_name_lower: {"employees": int, "source": str}}
_company_cache = {}

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = 64 * 1024 * 1024

# ── Paths ─────────────────────────────────────────────────────────────────────
BASE_DIR      = os.path.dirname(os.path.abspath(__file__))
SETTINGS_PATH = os.path.join(BASE_DIR, "settings.json")
SENT_LOG_PATH = os.path.join(BASE_DIR, "sent_log.json")

# Running on Render (or any cloud host) vs local Windows machine
IS_CLOUD = os.getenv("RENDER") == "true" or os.getenv("IS_CLOUD") == "true"

# ── Reference library dirs ────────────────────────────────────────────────────
# Three tiers — Claude reads all of them when generating emails.
#
#  overall/    → FEAAM company docs, tone/style guides, correction notes
#                Claude uses these as base context for EVERY email
#
#  followup1/  → Add-value content (sector decks, PPTX)
#                Claude generates Follow-up 1 emails around these + attaches them
#
#  followup2/  → Deeper content (case studies, technical papers, PDFs)
#                Claude generates Follow-up 2 emails around these + attaches them
#
# Files persist until manually replaced or deleted via the Library tab.

LIBRARY_DIR       = os.path.join(BASE_DIR, "library")
OVERALL_DIR       = os.path.join(LIBRARY_DIR, "overall")      # tier 1
FOLLOWUP1_DIR     = os.path.join(LIBRARY_DIR, "followup1")    # tier 2
FOLLOWUP2_DIR     = os.path.join(LIBRARY_DIR, "followup2")    # tier 3
LIBRARY_META_PATH = os.path.join(LIBRARY_DIR, "meta.json")

# Legacy aliases so existing send logic still works
DECK_DIR       = FOLLOWUP1_DIR
CASE_STUDY_DIR = FOLLOWUP2_DIR

for _d in (LIBRARY_DIR, OVERALL_DIR, FOLLOWUP1_DIR, FOLLOWUP2_DIR):
    os.makedirs(_d, exist_ok=True)

# On first run (local), symlink / copy existing local files into the library
def _seed_local_files():
    """Copy existing local files into the library tiers on first run."""
    import shutil

    # Follow-up 1 — sector decks
    local_deck_src = r"C:\Users\mypc\Downloads"
    for fname in SECTOR_DECKS.values():
        src = os.path.join(local_deck_src, fname)
        dst = os.path.join(FOLLOWUP1_DIR, fname)
        if os.path.exists(src) and not os.path.exists(dst):
            try: shutil.copy2(src, dst)
            except Exception: pass

    # Follow-up 2 — case studies
    local_cs_src = r"C:\Users\mypc\OneDrive\Desktop\New folder\CASE STUDIES"
    for _, fname in CASE_STUDIES:
        src = os.path.join(local_cs_src, fname)
        dst = os.path.join(FOLLOWUP2_DIR, fname)
        if os.path.exists(src) and not os.path.exists(dst):
            try: shutil.copy2(src, dst)
            except Exception: pass

    # Overall reference — seed known briefing docs
    overall_sources = [
        r"C:\Users\mypc\OneDrive\Desktop\CLAUDE\feaam-matcher\reference_files\FEAAM_Technology.pdf",
        r"C:\Users\mypc\Downloads\Future Growth Plan FEAAM (1).pptx",
        r"C:\Users\mypc\OneDrive\Desktop\New folder\Christian's Remark.pdf",
        r"C:\Users\mypc\OneDrive\Desktop\New folder\Ebru's Feedback (Per mail).docx",
    ]
    for src in overall_sources:
        if os.path.exists(src):
            dst = os.path.join(OVERALL_DIR, os.path.basename(src))
            if not os.path.exists(dst):
                try: shutil.copy2(src, dst)
                except Exception: pass

# ── Sector → deck filename mapping ────────────────────────────────────────────
SECTOR_DECKS = {
    "humanoid_robot":  "FEAAM_Humanoid_Robot_Drive_Strategy.pptx",
    "forklift":        "Forklift_EV_Drive_Strategy.pptx",
    "drone":           "FEAAM_Drone_Drive_Strategy.pptx",
    "2_3_wheeler":     "FEAAM_2_3_Wheeler_Drive_Strategy.pptx",
    "hvac":            "FEAAM_HVAC_Drive_Strategy.pptx",
    "pump":            "FEAAM_ElectricalSubmersiblePump_Drive_Strategy.pptx",
    "wind":            "FEAAM_Wind_Energy_Drive_Strategy.pptx",
    "gaming_ffb":      "FEAAM_FFB_Simulation_Drive_Strategy.pptx",
}

# Keywords → sector key
SECTOR_KEYWORDS = [
    (["humanoid", "robot", "robotics"],                           "humanoid_robot"),
    (["forklift", "lift truck", "material handling", "agv"],      "forklift"),
    (["drone", "uav", "uas", "unmanned aerial"],                  "drone"),
    (["2 wheel", "2-wheel", "3 wheel", "3-wheel", "e-bike",
      "ebike", "scooter", "moped", "motorbike", "motorcycle"],    "2_3_wheeler"),
    (["hvac", "heating", "ventilation", "air condition",
      "cooling system", "chiller"],                               "hvac"),
    (["pump", "submersible", "hydraulic pump",
      "fluid", "oil & gas pump"],                                 "pump"),
    (["wind energy", "wind turbine", "wind power",
      "wind farm", "offshore wind"],                              "wind"),
    (["gaming", "simulator", "simulation", "force feedback",
      "ffb", "steering wheel", "haptic"],                         "gaming_ffb"),
    (["aerospace", "aviation", "satellite", "spacecraft"],        "drone"),  # closest
]

CASE_STUDIES = [
    ("Case Study 1 – Single-Sided Axial Flux Motor",
     "RE - Case Study 1 - Bratke_single_sided_AFM_authors_V2.pdf"),
    ("Case Study 2 – Drive Control (EDPE 2025)",
     "RE - Case Study 2 - EDPE2025_Avci.pdf"),
    ("Case Study 3 – Flux Barrier Cooling",
     "RE - Case Study 3 - A Flux Barrier Cooling.pdf"),
    ("Case Study 4 – Flux-Barrier eBike Motor",
     "RE- Case Study 4 - Flux-Barrier eBike Motor.pdf"),
    ("Case Study 5 – New PMSM Rotor (HUPM)",
     "RE - Case Study 5 - HUPM_new pmsm rotor.pdf"),
    ("Case Study 6 – Drive Technology (ICEMS 2025)",
     "RE - Case Study 6 - ICEMS2025_Avci.pdf"),
    ("Case Study 7 – Self-Excited Synchronous Machine",
     "RE - Case Study 7 - Self-excited SM.pdf"),
    ("Case Study 8 – Small Electric Motors",
     "RE - Case Study 8 - Small Electric Motors.pdf"),
]

SECTOR_CASE_STUDY_MAP = {
    "humanoid_robot": [4, 6],   # indices into CASE_STUDIES
    "forklift":       [2, 6],
    "drone":          [0, 4],
    "2_3_wheeler":    [3, 4],
    "hvac":           [2, 7],
    "pump":           [2, 6],
    "wind":           [0, 2],
    "gaming_ffb":     [1, 5],
}

# ── Sector display labels & email templates ───────────────────────────────────
SECTOR_LABELS = {
    "humanoid_robot": "Humanoid Robotics",
    "forklift":       "Electric Forklifts / Material Handling",
    "drone":          "Drones / UAV / Aerospace",
    "2_3_wheeler":    "Electric 2 & 3 Wheelers",
    "hvac":           "HVAC Systems",
    "pump":           "Electric Pumps / Submersible",
    "wind":           "Wind Energy",
    "gaming_ffb":     "Gaming / Force Feedback / Simulation",
}

SECTOR_EMAIL_TEMPLATES = {
    "humanoid_robot": {
        "subject": "Electric motor architecture for {company}'s humanoid robot development",
        "intro_exec": (
            "Humanoid robotics places some of the most demanding requirements on electric motors: "
            "compact form factor, high torque density at low speed, minimal vibration, and "
            "supply-chain resilience across 20–100 motors per platform.\n\n"
            "FEAAM's patented stator flux barrier motor architecture addresses all of these directly — "
            "reducing rare-earth magnet mass while maintaining or improving torque density and "
            "efficiency under identical electrical and geometrical boundary conditions."
        ),
        "intro_tech": (
            "The per-joint motor requirements in humanoid platforms — high torque density, "
            "minimal cogging, and supply-chain resilience — align closely with what FEAAM's "
            "patented stator flux barrier motor architecture delivers.\n\n"
            "Our architecture reduces rare-earth magnet mass while maintaining or improving "
            "torque density under identical electrical and geometrical boundary conditions, "
            "and has been validated using advanced electromagnetic simulation tools developed in-house."
        ),
        "cta": "drive strategy for humanoid robot applications",
    },
    "forklift": {
        "subject": "Drive efficiency for {company}'s electric forklift portfolio",
        "intro_exec": (
            "Electric forklifts and material handling systems demand high torque at near-zero speed, "
            "robust thermal performance under continuous duty, and reduced BOM cost pressure.\n\n"
            "FEAAM's patented stator flux barrier motor architecture reduces rare-earth magnet mass "
            "while improving torque density and efficiency — directly addressing the cost and "
            "performance trade-offs in industrial drive systems."
        ),
        "intro_tech": (
            "The drive requirements in electric forklift and material handling applications — "
            "high torque at standstill, thermal robustness under intermittent loading, and "
            "reduced rare-earth dependency — are precisely where FEAAM's patented stator flux "
            "barrier motor architecture delivers measurable improvements.\n\n"
            "Our architecture has been validated using advanced electromagnetic simulation and "
            "optimisation tools developed in-house at FEAAM."
        ),
        "cta": "electric forklift drive strategy",
    },
    "drone": {
        "subject": "High power density motor architecture for {company}'s drone / aerial applications",
        "intro_exec": (
            "Drone and UAV propulsion systems require the highest possible power-to-weight ratio, "
            "precise torque control, and robust operation under variable load — while managing "
            "rare-earth supply chain exposure as production volumes scale.\n\n"
            "FEAAM's patented stator flux barrier motor architecture achieves higher torque density "
            "in a smaller package, reducing rare-earth magnet mass without compromising performance."
        ),
        "intro_tech": (
            "Power-to-weight ratio and torque linearity are critical for drone propulsion. "
            "FEAAM's patented stator flux barrier motor architecture reduces rare-earth magnet mass "
            "while maintaining or improving torque density under identical geometrical boundary "
            "conditions — validated using advanced electromagnetic simulation tools developed in-house."
        ),
        "cta": "drone and aerial propulsion drive strategy",
    },
    "2_3_wheeler": {
        "subject": "Motor technology for {company}'s electric 2/3-wheeler powertrain",
        "intro_exec": (
            "Electric 2 and 3-wheeler powertrains face intense cost pressure, rapidly growing volumes, "
            "and increasing rare-earth supply risk as the market scales exponentially.\n\n"
            "FEAAM's patented stator flux barrier motor architecture reduces rare-earth magnet mass "
            "while maintaining or improving torque density and efficiency — offering a direct route "
            "to BOM cost reduction and supply-chain de-risking at scale."
        ),
        "intro_tech": (
            "Reducing rare-earth magnet dependency at the motor architecture level is one of the most "
            "direct levers for cost reduction in electric 2/3-wheeler powertrains.\n\n"
            "FEAAM's patented stator flux barrier motor architecture reduces magnet mass while "
            "maintaining or improving torque density and efficiency under identical electrical and "
            "geometrical boundary conditions, supported by in-house electromagnetic simulation tools."
        ),
        "cta": "electric 2 & 3-wheeler drive strategy",
    },
    "hvac": {
        "subject": "Motor efficiency improvements for {company}'s HVAC applications",
        "intro_exec": (
            "HVAC systems operating under variable load demand high efficiency across a wide speed "
            "range, low noise, and long service life — with growing pressure to reduce energy "
            "consumption and rare-earth dependency in motor components.\n\n"
            "FEAAM's patented stator flux barrier motor architecture improves efficiency and "
            "reduces rare-earth magnet mass under identical electrical and geometrical boundary "
            "conditions — directly applicable to HVAC drive system design."
        ),
        "intro_tech": (
            "Efficiency at partial load, reduced cogging torque, and lower rare-earth content are "
            "the key architecture-level improvements available in HVAC motor design.\n\n"
            "FEAAM's patented stator flux barrier motor architecture delivers all three, validated "
            "using advanced electromagnetic simulation and optimisation tools developed in-house."
        ),
        "cta": "HVAC drive strategy",
    },
    "pump": {
        "subject": "Advanced motor architecture for {company}'s electric pump applications",
        "intro_exec": (
            "Electric submersible and industrial pump applications require sustained high torque, "
            "reliable thermal performance under continuous duty, and motors designed to operate "
            "in harsh environments with minimal maintenance exposure.\n\n"
            "FEAAM's patented stator flux barrier motor architecture improves torque density and "
            "efficiency while reducing rare-earth magnet mass — applicable to both surface and "
            "submersible pump drive systems."
        ),
        "intro_tech": (
            "Torque density at continuous duty and thermal robustness are the dominant architecture "
            "constraints in electric pump motor design.\n\n"
            "FEAAM's patented stator flux barrier motor architecture addresses both — reducing "
            "rare-earth magnet mass while maintaining or improving torque density under identical "
            "electrical and geometrical boundary conditions, with in-house simulation validation."
        ),
        "cta": "electrical submersible pump drive strategy",
    },
    "wind": {
        "subject": "Generator architecture for {company}'s wind energy systems",
        "intro_exec": (
            "Wind energy generators face long service life requirements, extreme rare-earth "
            "supply exposure at scale, and increasing demand for magnet-free alternatives "
            "as installed capacity grows.\n\n"
            "FEAAM offers both reduced-magnet (flux barrier stator) and fully magnet-free "
            "synchronous machine architectures — providing supply-chain independence while "
            "maintaining power density equivalent to conventional permanent magnet designs."
        ),
        "intro_tech": (
            "Reducing rare-earth dependency at the generator architecture level is a strategic "
            "priority for wind energy OEMs and developers.\n\n"
            "FEAAM's patented stator flux barrier architecture and magnet-free synchronous machine "
            "technology both address this directly, validated using advanced electromagnetic "
            "simulation tools developed in-house at FEAAM."
        ),
        "cta": "wind energy drive strategy",
    },
    "gaming_ffb": {
        "subject": "Direct-drive motor architecture for {company}'s simulation / force feedback products",
        "intro_exec": (
            "High-performance simulation and force feedback applications demand exceptional torque "
            "linearity, minimal cogging, high bandwidth, and precise position control — areas where "
            "motor architecture has a direct and measurable impact on product quality.\n\n"
            "FEAAM's patented stator flux barrier motor architecture delivers higher torque density "
            "and reduced cogging in a compact design, directly improving haptic fidelity and "
            "dynamic response in direct-drive simulator systems."
        ),
        "intro_tech": (
            "Cogging torque reduction and high torque bandwidth are the critical motor architecture "
            "parameters in force feedback and direct-drive simulation applications.\n\n"
            "FEAAM's patented stator flux barrier motor architecture addresses both, reducing "
            "rare-earth magnet mass while improving torque density and linearity — validated "
            "using advanced electromagnetic simulation tools developed in-house."
        ),
        "cta": "force feedback and simulation drive strategy",
    },
}

SIGNATURE = "Prof. Dr.-Ing. Dieter Gerling\nFounder, FEAAM GmbH"

TOP_N = 25          # How many top-scored leads to show after upload

# How many contacts to allow per company based on priority/size:
#   P1 (medium 250-5000):  up to 2 contacts — different roles worth targeting
#   P2 (large  5000+):     up to 3 contacts — multiple decision-makers
#   P0 (unknown):          1 contact only
MAX_CONTACTS_PER_COMPANY = {1: 2, 2: 3, 0: 1}

EXEC_KEYWORDS = {
    "president","ceo","cto","cfo","coo","cso","chief","founder","owner",
    "partner","chairman","board","executive","evp","svp","vice president",
    "director","head of","vp",
}

# ── Helpers ───────────────────────────────────────────────────────────────────

def safe_str(v):
    s = str(v).strip()
    return "" if s.lower() in ("nan", "none", "") else s


def is_exec(title):
    t = (title or "").lower()
    return any(k in t for k in EXEC_KEYWORDS)


def detect_sector(text):
    """Return sector key from any free-text field."""
    t = (text or "").lower()
    for keywords, sector in SECTOR_KEYWORDS:
        if any(kw in t for kw in keywords):
            return sector
    return None


# ── Built-in company size database (common industrial companies) ──────────────
# tier: "medium" = 250-5000, "large" = 5000+
_KNOWN_COMPANIES = {
    # Large (>5000)
    "volvo":          {"employees": 100000, "tier": "large",  "note": "Global truck/industrial OEM"},
    "volvo group":    {"employees": 100000, "tier": "large",  "note": "Global truck/industrial OEM"},
    "chevron":        {"employees": 45000,  "tier": "large",  "note": "Global energy company"},
    "nov":            {"employees": 34000,  "tier": "large",  "note": "Oilfield equipment OEM"},
    "solar turbines": {"employees": 8000,   "tier": "large",  "note": "Gas turbine manufacturer (Caterpillar subsidiary)"},
    "siemens":        {"employees": 320000, "tier": "large",  "note": "Global industrial conglomerate"},
    "bosch":          {"employees": 430000, "tier": "large",  "note": "Global automotive/industrial supplier"},
    "abb":            {"employees": 105000, "tier": "large",  "note": "Power and automation tech"},
    "danfoss":        {"employees": 42000,  "tier": "large",  "note": "Industrial drives & HVAC"},
    "nidec":          {"employees": 140000, "tier": "large",  "note": "Electric motor manufacturer"},
    "honda":          {"employees": 211000, "tier": "large",  "note": "Global automotive OEM"},
    "toyota":         {"employees": 375000, "tier": "large",  "note": "Global automotive OEM"},
    "bmw":            {"employees": 149000, "tier": "large",  "note": "Global automotive OEM"},
    "rivian":         {"employees": 14000,  "tier": "large",  "note": "EV manufacturer"},
    "xiaomi":         {"employees": 35000,  "tier": "large",  "note": "Consumer electronics / EV"},
    "ge":             {"employees": 170000, "tier": "large",  "note": "Industrial conglomerate"},
    "general electric":{"employees": 170000,"tier": "large",  "note": "Industrial conglomerate"},
    "vestas":         {"employees": 30000,  "tier": "large",  "note": "Wind turbine OEM"},
    "enercon":        {"employees": 22000,  "tier": "large",  "note": "Wind turbine OEM"},
    "grundfos":       {"employees": 20000,  "tier": "large",  "note": "Pump manufacturer"},
    "flowserve":      {"employees": 17000,  "tier": "large",  "note": "Pump and valve OEM"},
    "parker hannifin":{"employees": 58000,  "tier": "large",  "note": "Motion & control"},
    "emerson":        {"employees": 88000,  "tier": "large",  "note": "Automation technology"},
    "honeywell":      {"employees": 99000,  "tier": "large",  "note": "Industrial conglomerate"},
    "carrier":        {"employees": 55000,  "tier": "large",  "note": "HVAC systems OEM"},
    "trane":          {"employees": 40000,  "tier": "large",  "note": "HVAC systems"},
    "daikin":         {"employees": 98000,  "tier": "large",  "note": "HVAC manufacturer"},
    "dji":            {"employees": 14000,  "tier": "large",  "note": "Consumer/commercial drone OEM"},
    "agility robotics":{"employees": 1500,  "tier": "medium", "note": "Humanoid robot startup"},
    "boston dynamics": {"employees": 600,   "tier": "medium", "note": "Robotics company (Hyundai)"},
    "figure ai":      {"employees": 300,    "tier": "medium", "note": "Humanoid robot startup"},
    "1x technologies":{"employees": 200,    "tier": "small",  "note": "Humanoid robot startup"},
    "toyota industries":{"employees": 50000,"tier": "large",  "note": "Forklift / industrial OEM"},
    "jungheinrich":   {"employees": 20000,  "tier": "large",  "note": "Forklift OEM"},
    "kion":           {"employees": 43000,  "tier": "large",  "note": "Forklift / warehouse OEM"},
    "crown equipment":{"employees": 19000,  "tier": "large",  "note": "Forklift manufacturer"},
    "hyster-yale":    {"employees": 7000,   "tier": "large",  "note": "Forklift OEM"},
    "raymond":        {"employees": 3000,   "tier": "medium", "note": "Forklift manufacturer"},
    "logisnext":      {"employees": 10000,  "tier": "large",  "note": "Forklift OEM (Mitsubishi)"},
    "ola electric":   {"employees": 10000,  "tier": "large",  "note": "E-scooter OEM"},
    "hero electric":  {"employees": 2000,   "tier": "medium", "note": "E-2 wheeler OEM"},
    "bajaj":          {"employees": 21000,  "tier": "large",  "note": "2/3-wheeler OEM"},
    "tvs motor":      {"employees": 9000,   "tier": "large",  "note": "2/3-wheeler OEM"},
    "cpac":              {"employees": 500,    "tier": "medium", "note": "Industrial controls"},
    "leistritz":         {"employees": 3000,  "tier": "medium", "note": "Industrial pump/extrusion OEM"},
    # Oil & gas / industrial pump sector
    "nov":               {"employees": 34000, "tier": "large",  "note": "Oilfield equipment OEM", "sector_hint": "pump"},
    "national oilwell":  {"employees": 34000, "tier": "large",  "note": "Oilfield equipment OEM", "sector_hint": "pump"},
    "baker hughes":      {"employees": 54000, "tier": "large",  "note": "Oilfield services", "sector_hint": "pump"},
    "schlumberger":      {"employees": 98000, "tier": "large",  "note": "Oilfield services", "sector_hint": "pump"},
    "slb":               {"employees": 98000, "tier": "large",  "note": "Oilfield services (SLB)", "sector_hint": "pump"},
    "halliburton":       {"employees": 48000, "tier": "large",  "note": "Oilfield services", "sector_hint": "pump"},
    "weatherford":       {"employees": 17000, "tier": "large",  "note": "Oilfield equipment", "sector_hint": "pump"},
    "solar turbines":    {"employees": 8000,  "tier": "large",  "note": "Gas turbine manufacturer (Caterpillar)", "sector_hint": "pump"},
    "gardner denver":    {"employees": 7000,  "tier": "large",  "note": "Industrial compressors/pumps", "sector_hint": "pump"},
    "sulzer":            {"employees": 14000, "tier": "large",  "note": "Pump & rotating equipment OEM", "sector_hint": "pump"},
    "xylem":             {"employees": 21000, "tier": "large",  "note": "Water technology / pumps", "sector_hint": "pump"},
    "ebara":             {"employees": 13000, "tier": "large",  "note": "Pump manufacturer", "sector_hint": "pump"},
    "itt":               {"employees": 10000, "tier": "large",  "note": "Industrial pumps & motion", "sector_hint": "pump"},
    "gorman rupp":       {"employees": 1700,  "tier": "medium", "note": "Industrial pump OEM", "sector_hint": "pump"},
    "pentair":           {"employees": 11000, "tier": "large",  "note": "Water treatment / pumps", "sector_hint": "pump"},
    "wilo":              {"employees": 9000,  "tier": "large",  "note": "Pump manufacturer", "sector_hint": "pump"},
    "ksb":               {"employees": 15000, "tier": "large",  "note": "Pump and valve OEM", "sector_hint": "pump"},
    "chevron":           {"employees": 45000, "tier": "large",  "note": "Global energy company (oil & gas)", "sector_hint": "pump"},
    "exxon":             {"employees": 62000, "tier": "large",  "note": "Global energy company", "sector_hint": "pump"},
    "exxonmobil":        {"employees": 62000, "tier": "large",  "note": "Global energy company", "sector_hint": "pump"},
    "shell":             {"employees": 93000, "tier": "large",  "note": "Global energy company", "sector_hint": "pump"},
    "bp":                {"employees": 67000, "tier": "large",  "note": "Global energy company", "sector_hint": "pump"},
    "totalenergies":     {"employees": 101000,"tier": "large",  "note": "Global energy company", "sector_hint": "pump"},
    "magvar":            {"employees": 200,   "tier": "small",  "note": "Magnetic variance / navigation tech", "sector_hint": "forklift"},
}

def _lookup_known(company: str) -> dict | None:
    """Check built-in database (case-insensitive, fuzzy prefix match)."""
    key = company.strip().lower()
    if key in _KNOWN_COMPANIES:
        return {**_KNOWN_COMPANIES[key], "source": "database"}
    for known_key, data in _KNOWN_COMPANIES.items():
        if known_key in key or key.startswith(known_key[:6]):
            return {**data, "source": "database"}
    return None


# Built-in FEAAM fit descriptions per sector (used when Claude API not available)
_SECTOR_FEAAM_FIT = {
    "pump": (
        "Electric submersible and high-pressure pump drives require sustained high torque under continuous duty and reliable operation in harsh environments. "
        "FEAAM's patented stator flux barrier motor architecture improves torque density and efficiency while reducing rare-earth magnet mass — directly reducing BOM cost and supply-chain exposure for pump drive systems."
    ),
    "forklift": (
        "Electric forklift drives demand high torque at near-zero speed, thermal robustness under intermittent loading, and competitive BOM cost. "
        "FEAAM's patented stator flux barrier motor architecture delivers higher torque density under identical geometrical constraints, with reduced rare-earth magnet mass — a direct cost and performance advantage for material handling OEMs."
    ),
    "humanoid_robot": (
        "Humanoid robot platforms require 20–100 compact, high-torque-density motors per unit, with minimal rare-earth dependency at scale. "
        "FEAAM's patented stator flux barrier motor architecture achieves higher torque density in a smaller package — reducing rare-earth magnet mass without compromising performance, critical for viable humanoid robot economics."
    ),
    "drone": (
        "Drone and UAV propulsion requires maximum power-to-weight ratio and precise torque linearity in a compact form factor. "
        "FEAAM's stator flux barrier architecture delivers higher torque density while reducing rare-earth magnet mass — directly improving flight time and reducing supply-chain risk as production volumes scale."
    ),
    "2_3_wheeler": (
        "Electric 2 and 3-wheeler powertrains face intense cost pressure and growing rare-earth supply risk as volumes scale exponentially. "
        "FEAAM's patented stator flux barrier motor architecture reduces rare-earth magnet mass while maintaining torque density — a direct BOM cost reduction lever for volume e-mobility manufacturers."
    ),
    "hvac": (
        "HVAC drive systems require high efficiency across a wide speed range and long service life with low maintenance exposure. "
        "FEAAM's stator flux barrier motor architecture improves part-load efficiency and reduces cogging torque — translating to lower energy consumption and quieter operation in HVAC applications."
    ),
    "wind": (
        "Wind energy generators face significant rare-earth supply exposure at scale and increasing demand for magnet-free alternatives. "
        "FEAAM offers both reduced-magnet flux barrier and fully magnet-free synchronous machine architectures — providing supply-chain independence while maintaining power density equivalent to conventional PM designs."
    ),
    "gaming_ffb": (
        "Force feedback and direct-drive simulation systems require exceptional torque linearity, minimal cogging torque, and high bandwidth. "
        "FEAAM's stator flux barrier motor architecture reduces cogging and improves torque density in a compact design — directly improving haptic fidelity and dynamic response for professional simulation products."
    ),
}


def _lookup_wikipedia(company: str) -> dict | None:
    """Query Wikipedia infobox for employee count. No API key needed."""
    import urllib.request as _ureq, urllib.parse as _uparse
    try:
        # Search for best matching page
        search_url = "https://en.wikipedia.org/w/api.php?" + _uparse.urlencode({
            "action": "query", "list": "search",
            "srsearch": company + " company",
            "format": "json", "srlimit": 2,
        })
        req = _ureq.Request(search_url, headers={"User-Agent": "FEAAM-Sender/1.0"})
        with _ureq.urlopen(req, timeout=6) as r:
            hits = json.loads(r.read()).get("query", {}).get("search", [])
        if not hits:
            return None
        title = hits[0]["title"]

        # Get page wikitext (section 0 = intro + infobox)
        page_url = "https://en.wikipedia.org/w/api.php?" + _uparse.urlencode({
            "action": "query", "titles": title, "prop": "revisions",
            "rvprop": "content", "rvslots": "main",
            "format": "json", "rvsection": "0",
        })
        req2 = _ureq.Request(page_url, headers={"User-Agent": "FEAAM-Sender/1.0"})
        with _ureq.urlopen(req2, timeout=6) as r:
            pages = json.loads(r.read()).get("query", {}).get("pages", {})
        wikitext = (list(pages.values())[0]
                    .get("revisions", [{}])[0]
                    .get("slots", {}).get("main", {}).get("*", ""))

        # Extract employee count from infobox
        m = re.search(r"num_employees\s*=\s*([\d,]+)", wikitext)
        if m:
            emp = int(m.group(1).replace(",", ""))
            tier = "medium" if 250 <= emp <= 5000 else ("large" if emp > 5000 else "small")
            return {"employees": emp, "tier": tier,
                    "note": f"Wikipedia: {title}", "source": "wikipedia"}
    except Exception:
        pass
    return None


def _lookup_wikipedia_full(company: str) -> dict:
    """
    Fetch Wikipedia intro text and infobox for a company.
    Returns dict with employees, description, products fields.
    """
    import urllib.request as _ureq, urllib.parse as _uparse
    result = {}
    try:
        search_url = "https://en.wikipedia.org/w/api.php?" + _uparse.urlencode({
            "action": "query", "list": "search",
            "srsearch": company + " company",
            "format": "json", "srlimit": 2,
        })
        req = _ureq.Request(search_url, headers={"User-Agent": "FEAAM-Sender/1.0"})
        with _ureq.urlopen(req, timeout=8) as r:
            hits = json.loads(r.read()).get("query", {}).get("search", [])
        if not hits:
            return result
        title = hits[0]["title"]

        # Get intro text (plain extract)
        extract_url = "https://en.wikipedia.org/w/api.php?" + _uparse.urlencode({
            "action": "query", "titles": title, "prop": "extracts",
            "exintro": True, "explaintext": True,
            "format": "json",
        })
        req3 = _ureq.Request(extract_url, headers={"User-Agent": "FEAAM-Sender/1.0"})
        with _ureq.urlopen(req3, timeout=8) as r:
            pages = json.loads(r.read()).get("query", {}).get("pages", {})
        extract = (list(pages.values())[0].get("extract", "") or "")[:2000]

        # Get wikitext for infobox employee count
        page_url = "https://en.wikipedia.org/w/api.php?" + _uparse.urlencode({
            "action": "query", "titles": title, "prop": "revisions",
            "rvprop": "content", "rvslots": "main",
            "format": "json", "rvsection": "0",
        })
        req2 = _ureq.Request(page_url, headers={"User-Agent": "FEAAM-Sender/1.0"})
        with _ureq.urlopen(req2, timeout=8) as r:
            pages2 = json.loads(r.read()).get("query", {}).get("pages", {})
        wikitext = (list(pages2.values())[0]
                    .get("revisions", [{}])[0]
                    .get("slots", {}).get("main", {}).get("*", ""))

        # Extract employee count
        m = re.search(r"num_employees\s*=\s*([\d,]+)", wikitext)
        if m:
            emp = int(m.group(1).replace(",", ""))
            tier = "medium" if 250 <= emp <= 5000 else ("large" if emp > 5000 else "small")
            result["employees"] = emp
            result["tier"]      = tier

        result["description"] = extract.strip()
        result["wiki_title"]  = title
        result["source"]      = "wikipedia"
    except Exception as e:
        print(f"[wiki_full] {company}: {e}")
    return result


def deep_research_company(company: str, contact_name: str = "", contact_title: str = "",
                           sector: str = "", location: str = "") -> dict:
    """
    Deep-research one company using Wikipedia + Claude (if available).
    Returns:
      {
        "employees": int,
        "tier": str,
        "description": str,
        "products": str,
        "feaam_fit": str,
        "research_note": str,
        "source": str,
      }
    """
    result = {
        "employees": 0, "tier": "unknown",
        "description": "", "products": "",
        "feaam_fit": "", "research_note": "",
        "source": "none",
    }

    # ── Pass 1: built-in database ─────────────────────────────────
    db_hit = _lookup_known(company)
    if db_hit:
        result.update(db_hit)

    # ── Pass 2: Wikipedia ─────────────────────────────────────────
    wiki = _lookup_wikipedia_full(company)
    if wiki:
        if wiki.get("employees") and not result.get("employees"):
            result["employees"] = wiki["employees"]
            result["tier"]      = wiki["tier"]
        if wiki.get("description"):
            result["description"] = wiki["description"]
        result["source"] = "wikipedia"

    # ── Pass 3: Claude API — synthesis & FEAAM fit analysis ───────
    ai, err = _get_ai_client()
    if ai:
        sector_label = SECTOR_LABELS.get(sector, sector or "industrial electric motors")
        wiki_context = result["description"][:800] if result["description"] else "No Wikipedia data found."
        prompt = f"""You are a B2B sales research assistant for FEAAM GmbH, a German electric motor
technology company founded by Prof. Dr.-Ing. Dieter Gerling (former chair of electrical drives at
Bundeswehr University Munich). FEAAM's core IP is a patented stator flux barrier motor architecture
that reduces rare-earth magnet mass while maintaining or improving torque density and efficiency.
FEAAM also offers magnet-free synchronous machine technology. They serve: humanoid robotics,
electric forklifts, drones/UAV, e-bikes/scooters, HVAC, electric pumps, wind energy, and
gaming/force-feedback simulation.

Research target:
- Company: {company}
- Contact: {contact_name or "unknown"} ({contact_title or "unknown role"})
- Sector hint: {sector_label}
- Location: {location or "unknown"}

Wikipedia summary (may be incomplete or wrong company):
{wiki_context}

Based on your training knowledge about {company}, provide a JSON object with these fields:
{{
  "employees": <integer estimate, 0 if truly unknown>,
  "tier": "<small|medium|large|unknown>  (small<250, medium 250-5000, large>5000)",
  "description": "<2-3 sentence factual description of what {company} actually makes/does>",
  "products": "<comma-separated list of {company}'s key products or application areas most relevant to electric motors>",
  "feaam_fit": "<2-3 sentences explaining specifically why FEAAM's motor technology would benefit {company}'s products — be concrete, reference their actual products>",
  "research_note": "<one sentence on {contact_name}'s likely decision-making role or how to personalise the pitch to a {contact_title}>"
}}

Return ONLY valid JSON. Be factual — if you don't know the company well, say so briefly in description."""

        try:
            resp = ai.messages.create(
                model="claude-haiku-4-5", max_tokens=600,
                messages=[{"role": "user", "content": prompt}]
            )
            raw  = re.sub(r"^```(?:json)?\s*|\s*```$", "", resp.content[0].text.strip())
            data = json.loads(raw)
            # Merge — Claude overrides Wikipedia where it has data
            if data.get("employees"): result["employees"] = data["employees"]
            if data.get("tier") and data["tier"] != "unknown": result["tier"] = data["tier"]
            for k in ("description", "products", "feaam_fit", "research_note"):
                if data.get(k): result[k] = data[k]
            result["source"] = "ai"
        except Exception as e:
            print(f"[deep_research] Claude error for {company}: {e}")

    # ── Pass 4: Built-in FEAAM fit text (fallback when no Claude API) ────────
    if not result.get("feaam_fit"):
        # Use sector hint from DB, or the passed sector argument
        effective_sector = result.get("sector_hint") or sector or ""
        fit_text = _SECTOR_FEAAM_FIT.get(effective_sector, "")
        if not fit_text and result.get("description"):
            # Generic fit based on description
            fit_text = (
                f"FEAAM's patented stator flux barrier motor architecture — which reduces "
                f"rare-earth magnet mass while improving torque density and efficiency — "
                f"is applicable to {company}'s electric drive systems, offering BOM cost "
                f"reduction and supply-chain de-risking."
            )
        result["feaam_fit"] = fit_text

    # ── Finalise tier from employee count ─────────────────────────────────────
    if result["tier"] == "unknown" and result["employees"]:
        emp = result["employees"]
        result["tier"] = "medium" if 250 <= emp <= 5000 else ("large" if emp > 5000 else "small")

    return result


def get_company_size(company_name: str) -> dict:
    """Return cached size info for a company, or empty dict."""
    key = company_name.strip().lower()
    return _company_cache.get(key, {})


def score_lead(row_dict):
    """
    Returns (priority, sector_key, score_breakdown)
    priority: 1, 2, or 0 (unknown)
    """
    # --- Employee count: check spreadsheet columns first, then AI cache ---
    emp = 0
    emp_source = "sheet"
    for k in ("employees", "employee_count", "headcount", "staff", "size", "num_employees"):
        v = row_dict.get(k) or row_dict.get(k.replace("_", " ")) or row_dict.get(k.title())
        if v:
            try:
                emp = int(str(v).replace(",", "").split("-")[0].strip())
                break
            except Exception:
                pass

    # Fall back to AI research cache
    if emp == 0:
        company = row_dict.get("company", "")
        cached = get_company_size(company)
        if cached.get("employees"):
            emp = cached["employees"]
            emp_source = "ai"

    if 250 <= emp <= 5000:
        priority = 1
    elif emp > 5000:
        priority = 2
    else:
        priority = 0  # unknown

    # --- Sector ---
    sector = None
    # First check built-in DB for sector_hint (most reliable)
    db = _lookup_known(row_dict.get("company", ""))
    if db and db.get("sector_hint"):
        sector = db["sector_hint"]

    if not sector:
        for k in ("matched_sector", "sector", "industry", "business_type", "segment"):
            v = row_dict.get(k) or row_dict.get(k.replace("_", " ")) or row_dict.get(k.title())
            if v:
                s = detect_sector(str(v))
                if s:
                    sector = s
                    break
    if not sector:
        combined = " ".join([
            str(row_dict.get("company", "")),
            str(row_dict.get("job_title", "")),
            str(row_dict.get("jobtitle", "")),
        ])
        sector = detect_sector(combined)

    # --- Numeric score (higher = better) ---
    prio_score = {1: 100, 2: 60, 0: 20}[priority]
    sector_score = 20 if sector else 0
    classification = str(row_dict.get("classification", "")).upper()
    fit_score = 15 if classification == "FIT" else (0 if classification == "UNFIT" else 5)
    try:
        conf_score = float(str(row_dict.get("confidence", 50) or 50).replace("%", "")) / 10
    except Exception:
        conf_score = 5

    total = prio_score + sector_score + fit_score + conf_score

    return priority, sector, total


def generate_email(lead):
    """
    Build subject + body for an add-value content email.
    Uses deep-research data (feaam_fit, products, description) when available
    to personalise the email beyond the generic sector template.
    """
    sector       = lead.get("sector")
    company      = lead.get("company", "your company")
    first        = lead.get("first_name") or (lead.get("full_name") or "").split()[0] or "there"
    title        = lead.get("job_title", "")
    level        = "exec" if is_exec(title) else "tech"
    feaam_fit    = (lead.get("feaam_fit") or "").strip()
    products     = (lead.get("products") or "").strip()
    research_note= (lead.get("research_note") or "").strip()

    tmpl = SECTOR_EMAIL_TEMPLATES.get(sector, SECTOR_EMAIL_TEMPLATES["humanoid_robot"])
    cta_name     = tmpl["cta"]
    sector_label = SECTOR_LABELS.get(sector, "your sector")

    # ── Subject: use research-specific subject when products are known ─────
    if products:
        first_product = products.split(",")[0].strip()
        subject = f"FEAAM motor architecture for {company}'s {first_product} drive systems"
    else:
        subject = tmpl["subject"].format(company=company)

    # ── Intro paragraph ────────────────────────────────────────────────────
    if feaam_fit:
        # Research-driven personalised intro
        intro = feaam_fit
    else:
        intro = tmpl["intro_exec"] if level == "exec" else tmpl["intro_tech"]

    # ── Products mention (only if researched) ─────────────────────────────
    products_line = ""
    if products:
        products_line = (
            f"Your portfolio — including {products} — represents exactly the kind of "
            f"application where motor architecture choices have a direct impact on system "
            f"performance, BOM cost, and supply-chain resilience.\n\n"
        )

    # ── Personalisation note for contact role ─────────────────────────────
    role_line = ""
    if research_note and title:
        role_line = f"{research_note}\n\n"

    body = (
        f"Hi {first},\n\n"
        f"{intro}\n\n"
        f"{products_line}"
        f"{role_line}"
        f"I am attaching our {cta_name} deck for your reference — it outlines how FEAAM's "
        f"technology applies to {sector_label} applications, including specific performance "
        f"benchmarks and design considerations.\n\n"
        f"FEAAM is a spin-off of Bundeswehr University Munich with over 80 patents worldwide. "
        f"Our architecture has been evaluated across multiple industrial applications and is "
        f"supported by advanced electromagnetic simulation and optimisation tools developed in-house.\n\n"
        f"If the enclosed material is of interest, I would be glad to discuss how it applies "
        f"specifically to {company}'s requirements.\n\n"
        f"Best regards,\n{SIGNATURE}"
    )
    return subject, body


def get_deck_path(sector):
    fname = SECTOR_DECKS.get(sector, "")
    if not fname:
        return None, None
    path = os.path.join(DECK_DIR, fname)
    return (path if os.path.exists(path) else None), fname


def get_case_study_path(filename):
    path = os.path.join(CASE_STUDY_DIR, filename)
    return path if os.path.exists(path) else None


# ── Library metadata helpers ──────────────────────────────────────────────────

def load_library_meta():
    if os.path.exists(LIBRARY_META_PATH):
        with open(LIBRARY_META_PATH) as f:
            return json.load(f)
    return {"overall": [], "followup1": [], "followup2": []}


def save_library_meta(meta):
    with open(LIBRARY_META_PATH, "w") as f:
        json.dump(meta, f, indent=2)


def list_tier_files(tier_dir, meta_list):
    """Return list of file info dicts for a tier directory."""
    files = []
    # Files from meta (user-uploaded with labels)
    meta_fnames = {m["filename"] for m in meta_list}
    for m in meta_list:
        path = os.path.join(tier_dir, m["filename"])
        info = file_info(path, m["filename"])
        info["label"]       = m.get("label", m["filename"])
        info["description"] = m.get("description", "")
        info["uploaded_at"] = m.get("uploaded_at", "")
        files.append(info)
    # Also pick up any files present on disk that aren't in meta yet
    if os.path.isdir(tier_dir):
        for fname in sorted(os.listdir(tier_dir)):
            if fname not in meta_fnames and not fname.startswith("."):
                path = os.path.join(tier_dir, fname)
                if os.path.isfile(path):
                    info = file_info(path, fname)
                    info["label"]       = fname
                    info["description"] = ""
                    info["uploaded_at"] = info["modified"]
                    files.append(info)
    return files


def _read_text_from_file(path):
    """Extract plain text from PDF, DOCX, or PPTX for Claude context."""
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".pdf":
            import fitz
            doc  = fitz.open(path)
            text = "\n".join(p.get_text() for p in doc)
            doc.close()
            return text[:8000]
        elif ext == ".docx":
            from docx import Document
            doc = Document(path)
            return "\n".join(p.text for p in doc.paragraphs)[:8000]
        elif ext == ".pptx":
            from pptx import Presentation
            prs  = Presentation(path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text.append(shape.text.strip())
            return "\n".join(text)[:8000]
    except Exception as e:
        return f"[Could not read {os.path.basename(path)}: {e}]"
    return ""


def build_reference_context(tier="overall"):
    """
    Read all files in a library tier and return combined text for Claude.
    tier: 'overall' | 'followup1' | 'followup2'
    """
    dirs  = {"overall": OVERALL_DIR, "followup1": FOLLOWUP1_DIR, "followup2": FOLLOWUP2_DIR}
    folder = dirs.get(tier, OVERALL_DIR)
    if not os.path.isdir(folder):
        return ""
    parts = []
    for fname in sorted(os.listdir(folder)):
        if fname.startswith("."):
            continue
        path = os.path.join(folder, fname)
        if os.path.isfile(path):
            text = _read_text_from_file(path)
            if text:
                parts.append(f"=== {fname} ===\n{text}")
    return "\n\n".join(parts)


def file_info(path, fname):
    """Return size + modified date for a file."""
    try:
        st = os.stat(path)
        return {
            "filename": fname,
            "size_kb":  round(st.st_size / 1024, 1),
            "modified": datetime.fromtimestamp(st.st_mtime).strftime("%Y-%m-%d %H:%M"),
            "exists":   True,
        }
    except Exception:
        return {"filename": fname, "size_kb": 0, "modified": "", "exists": False}


# ── Settings ──────────────────────────────────────────────────────────────────

def load_settings():
    # Environment variables take priority (for Render / cloud deployments)
    env_user = os.getenv("SMTP_USER", "")
    env_pass = os.getenv("SMTP_PASS", "")
    if os.path.exists(SETTINGS_PATH):
        with open(SETTINGS_PATH) as f:
            s = json.load(f)
        if env_user: s["smtp_user"] = env_user
        if env_pass: s["smtp_pass"] = env_pass
        return s
    return {
        "smtp_host":        os.getenv("SMTP_HOST", "smtp.ionos.com"),
        "smtp_port":        int(os.getenv("SMTP_PORT", "587")),
        "smtp_user":        env_user,
        "smtp_pass":        env_pass,
        "sender_name":      os.getenv("SENDER_NAME",
                                      "Prof. Dr.-Ing. Dieter Gerling | FEAAM GmbH"),
        "anthropic_api_key": os.getenv("ANTHROPIC_API_KEY", ""),
    }


def save_settings(data):
    with open(SETTINGS_PATH, "w") as f:
        json.dump(data, f, indent=2)


# ── Sent log ──────────────────────────────────────────────────────────────────

def load_sent_log():
    if os.path.exists(SENT_LOG_PATH):
        with open(SENT_LOG_PATH) as f:
            return json.load(f)
    return []


def append_sent_log(entry):
    log = load_sent_log()
    log.append(entry)
    with open(SENT_LOG_PATH, "w") as f:
        json.dump(log, f, indent=2)


# ── In-memory session store ───────────────────────────────────────────────────
# Holds the last uploaded + prioritised leads list
_session = {"leads": [], "df": None, "research_status": "idle",
            "research_total": 0, "research_done": 0, "research_error": ""}


def normalise_columns(df):
    """Lowercase + strip column names, remove spaces/underscores for lookup."""
    df.columns = [str(c).strip() for c in df.columns]
    return df


def row_to_dict(row, df):
    """Map any column name variants to a clean dict."""
    d = {}
    for col in df.columns:
        key = col.lower().replace(" ", "_").replace("-", "_")
        d[key] = safe_str(row.get(col, ""))
    # Aliases
    d["full_name"]  = d.get("full_name") or f"{d.get('first_name','')} {d.get('last_name','')}".strip()
    d["first_name"] = d.get("first_name") or (d["full_name"].split()[0] if d["full_name"] else "")
    d["company"]    = d.get("company") or d.get("organization") or d.get("org") or ""
    d["job_title"]  = d.get("job_title") or d.get("jobtitle") or d.get("title") or d.get("position") or ""
    d["email"]      = d.get("email") or d.get("e_mail") or d.get("mail") or ""
    d["location"]   = d.get("location") or d.get("city") or d.get("country") or ""
    return d


# ── Routes ────────────────────────────────────────────────────────────────────

@app.route("/")
def index():
    return render_template("index.html")


def _build_leads(df, preserve_ids=None):
    """
    Score ALL leads, sort, return top TOP_N.
    preserve_ids: optional dict {company+email -> lead_id} to keep stable IDs on rebuild.
    """
    leads = []
    for _, row in df.iterrows():
        d = row_to_dict(row, df)
        if not d["email"] and not d["company"]:
            continue
        priority, sector, score = score_lead(d)

        # Preserve stable ID across rebuilds (so approval queue refs survive)
        stable_key = (d["company"] + "|" + d["email"]).lower()
        lead_id = (preserve_ids or {}).get(stable_key, str(uuid.uuid4()))

        subject, body = generate_email({**d, "sector": sector})
        deck_path, deck_fname = get_deck_path(sector)
        suggested_cs = [
            {"label": CASE_STUDIES[i][0], "filename": CASE_STUDIES[i][1]}
            for i in SECTOR_CASE_STUDY_MAP.get(sector, [])
        ]
        cached = get_company_size(d["company"])
        leads.append({
            "id":              lead_id,
            "company":         d["company"],
            "full_name":       d["full_name"],
            "first_name":      d["first_name"],
            "job_title":       d["job_title"],
            "email":           d["email"],
            "location":        d["location"],
            "sector":          sector,
            "sector_label":    SECTOR_LABELS.get(sector, "Unknown"),
            "priority":        priority,
            "score":           round(score, 1),
            "subject":         subject,
            "body":            body,
            "deck_fname":      deck_fname,
            "deck_exists":     deck_path is not None,
            "suggested_cs":    suggested_cs,
            "sent":            False,
            "employees":       cached.get("employees"),
            "emp_tier":        cached.get("tier", ""),
            "emp_note":        cached.get("note", ""),
            "emp_source":      cached.get("source", ""),
            "description":     cached.get("description", ""),
            "products":        cached.get("products", ""),
            "feaam_fit":       cached.get("feaam_fit", ""),
            "research_note":   cached.get("research_note", ""),
            "research_status": cached.get("_research_status", "pending"),
        })
    priority_order = {1: 0, 2: 1, 0: 2}
    leads.sort(key=lambda x: (priority_order[x["priority"]], -x["score"]))

    # Smart dedup — allow multiple contacts for important companies
    # P2 (large): up to 3 contacts | P1 (medium): up to 2 | Unknown: 1
    company_counts = {}
    deduped = []
    for lead in leads:
        co_key  = lead["company"].strip().lower()
        limit   = MAX_CONTACTS_PER_COMPANY.get(lead["priority"], 1)
        current = company_counts.get(co_key, 0)
        if current < limit:
            company_counts[co_key] = current + 1
            deduped.append(lead)
    leads = deduped

    return leads[:TOP_N]


@app.route("/api/upload", methods=["POST"])
def upload():
    f = request.files.get("file")
    if not f or not f.filename:
        return jsonify({"ok": False, "error": "No file uploaded"}), 400

    try:
        name = f.filename.lower()
        if name.endswith(".csv"):
            try:
                df = pd.read_csv(f)
            except Exception:
                f.seek(0)
                df = pd.read_csv(f, encoding="latin-1")
        else:
            df = pd.read_excel(f, engine="openpyxl")
        df = normalise_columns(df)
    except Exception as e:
        return jsonify({"ok": False, "error": str(e)}), 400

    # --- Phase 1: quick score without AI — return top TOP_N instantly ---
    leads = _build_leads(df)
    # Mark all as pending deep research
    for lead in leads:
        lead["research_status"] = "pending"

    _session["leads"]  = leads
    _session["df"]     = df
    _session["research_status"] = "running"
    _session["research_total"]  = len(leads)
    _session["research_done"]   = 0
    _session["research_error"]  = ""

    p1 = sum(1 for l in leads if l["priority"] == 1)
    p2 = sum(1 for l in leads if l["priority"] == 2)
    pu = sum(1 for l in leads if l["priority"] == 0)

    # --- Phase 2: deep research each of the top TOP_N one by one ---
    def _research_bg():
        try:
            for lead in _session["leads"]:
                if lead.get("research_status") == "skipped":
                    continue
                lead["research_status"] = "researching"
                company = lead.get("company", "")
                if not company:
                    lead["research_status"] = "skipped"
                    _session["research_done"] += 1
                    continue

                print(f"[research] Deep-researching: {company}")
                info = deep_research_company(
                    company=company,
                    contact_name=lead.get("full_name", ""),
                    contact_title=lead.get("job_title", ""),
                    sector=lead.get("sector", ""),
                    location=lead.get("location", ""),
                )

                # Update cache
                cache_key = company.strip().lower()
                _company_cache[cache_key] = info

                # Update lead in place with research results
                if info.get("employees"):
                    lead["employees"] = info["employees"]
                if info.get("tier") and info["tier"] != "unknown":
                    lead["emp_tier"] = info["tier"]
                    # Recalculate priority
                    emp = info["employees"]
                    if 250 <= emp <= 5000:
                        lead["priority"] = 1
                    elif emp > 5000:
                        lead["priority"] = 2
                for k in ("description", "products", "feaam_fit", "research_note"):
                    if info.get(k):
                        lead[k] = info[k]
                lead["emp_source"] = info.get("source", "")

                # If DB gave a sector hint and current sector is unknown, update it
                sector_hint = info.get("sector_hint")
                if sector_hint and (not lead.get("sector")):
                    lead["sector"]       = sector_hint
                    lead["sector_label"] = SECTOR_LABELS.get(sector_hint, sector_hint)
                    # Update attachment to match new sector
                    deck_path, deck_fname = get_deck_path(sector_hint)
                    lead["deck_fname"]  = deck_fname
                    lead["deck_exists"] = deck_path is not None
                    lead["suggested_cs"] = [
                        {"label": CASE_STUDIES[i][0], "filename": CASE_STUDIES[i][1]}
                        for i in SECTOR_CASE_STUDY_MAP.get(sector_hint, [])
                    ]

                # Regenerate personalised email with all research data
                new_subj, new_body = generate_email(lead)
                lead["subject"] = new_subj
                lead["body"]    = new_body

                lead["research_status"] = "done"
                _session["research_done"] += 1
                print(f"[research] Done: {company} ({_session['research_done']}/{_session['research_total']})")

            _session["research_status"] = "done"
        except Exception as e:
            print(f"[research_bg] {e}")
            _session["research_status"] = "error"
            _session["research_error"]  = str(e)

    threading.Thread(target=_research_bg, daemon=True).start()

    return jsonify({
        "ok":              True,
        "total":           len(leads),
        "p1":              p1,
        "p2":              p2,
        "unknown":         pu,
        "leads":           leads,
        "ai_researching":  True,
        "ai_count":        len(leads),
    })


@app.route("/api/research-status", methods=["GET"])
def research_status():
    """
    Frontend polls this throughout the research process.
    Returns per-lead status + full leads list on every poll so the UI
    can update each card as research completes.
    """
    status = _session.get("research_status", "idle")
    return jsonify({
        "ok":     True,
        "status": status,
        "total":  _session.get("research_total", 0),
        "done":   _session.get("research_done", 0),
        "error":  _session.get("research_error", ""),
        "leads":  _session["leads"],   # always return full list so UI updates progressively
    })


@app.route("/api/leads", methods=["GET"])
def get_leads():
    return jsonify({"ok": True, "leads": _session["leads"]})


@app.route("/api/lead/<lead_id>/email", methods=["POST"])
def update_email(lead_id):
    """Save edited subject/body back to the lead."""
    data = request.get_json()
    for lead in _session["leads"]:
        if lead["id"] == lead_id:
            lead["subject"] = data.get("subject", lead["subject"])
            lead["body"]    = data.get("body",    lead["body"])
            return jsonify({"ok": True})
    return jsonify({"ok": False, "error": "Lead not found"}), 404


@app.route("/api/lead/<lead_id>/regenerate", methods=["POST"])
def regenerate_email(lead_id):
    """Regenerate email for a lead from scratch."""
    for lead in _session["leads"]:
        if lead["id"] == lead_id:
            subject, body = generate_email(lead)
            lead["subject"] = subject
            lead["body"]    = body
            return jsonify({"ok": True, "subject": subject, "body": body})
    return jsonify({"ok": False, "error": "Lead not found"}), 404


@app.route("/api/send", methods=["POST"])
def send_email():
    """
    Send a single email with deck attachment.
    Requires explicit confirmation from the UI — this route will NOT be called
    unless the user clicked the final 'Confirm & Send' button.
    """
    data      = request.get_json()
    lead_id   = data.get("lead_id")
    subject   = data.get("subject", "")
    body      = data.get("body", "")
    cs_files  = data.get("case_studies", [])   # list of filenames to attach

    # Find lead
    lead = next((l for l in _session["leads"] if l["id"] == lead_id), None)
    if not lead:
        return jsonify({"ok": False, "error": "Lead not found"}), 404

    to_email = lead["email"]
    if not to_email:
        return jsonify({"ok": False, "error": "No email address for this lead"}), 400

    settings = load_settings()
    if not settings.get("smtp_user") or not settings.get("smtp_pass"):
        return jsonify({
            "ok": False,
            "error": "SMTP not configured. Go to Settings and enter your Ionos email and password."
        }), 400

    # Build email
    msg = EmailMessage()
    msg["Subject"] = subject
    msg["From"]    = f"{settings['sender_name']} <{settings['smtp_user']}>"
    msg["To"]      = to_email
    msg.set_content(body)

    # Attach sector deck (primary)
    if lead["sector"]:
        deck_path, _ = get_deck_path(lead["sector"])
        if deck_path and os.path.exists(deck_path):
            with open(deck_path, "rb") as fp:
                msg.add_attachment(
                    fp.read(),
                    maintype="application",
                    subtype="vnd.openxmlformats-officedocument.presentationml.presentation",
                    filename=os.path.basename(deck_path),
                )

    # Attach case studies (secondary, if selected)
    for cs_fname in cs_files:
        cs_path = get_case_study_path(cs_fname)
        if cs_path and os.path.exists(cs_path):
            with open(cs_path, "rb") as fp:
                msg.add_attachment(
                    fp.read(),
                    maintype="application",
                    subtype="pdf",
                    filename=cs_fname,
                )

    # Send
    try:
        with smtplib.SMTP(settings["smtp_host"], int(settings["smtp_port"])) as smtp:
            smtp.starttls()
            smtp.login(settings["smtp_user"], settings["smtp_pass"])
            smtp.send_message(msg)
    except Exception as e:
        return jsonify({"ok": False, "error": f"SMTP error: {str(e)}"}), 500

    # Mark as sent
    lead["sent"] = True
    lead["sent_at"] = datetime.now().isoformat(timespec="seconds")
    append_sent_log({
        "sent_at":   lead["sent_at"],
        "company":   lead["company"],
        "to_name":   lead["full_name"],
        "to_email":  to_email,
        "subject":   subject,
        "sector":    lead.get("sector_label"),
        "deck":      lead.get("deck_fname"),
        "case_studies": cs_files,
    })
    return jsonify({"ok": True, "message": f"Email sent to {to_email}"})


@app.route("/api/settings", methods=["GET"])
def get_settings():
    s = load_settings()
    s_safe = {k: v for k, v in s.items() if k not in ("smtp_pass", "anthropic_api_key")}
    s_safe["smtp_pass"]        = "••••••••" if s.get("smtp_pass")        else ""
    s_safe["anthropic_api_key"] = "••••••••" if s.get("anthropic_api_key") else ""
    return jsonify({"ok": True, "settings": s_safe})


@app.route("/api/settings", methods=["POST"])
def post_settings():
    data = request.get_json()
    current = load_settings()
    for key in ("smtp_host", "smtp_port", "smtp_user", "sender_name"):
        if key in data:
            current[key] = data[key]
    if data.get("smtp_pass") and not data["smtp_pass"].startswith("•"):
        current["smtp_pass"] = data["smtp_pass"]
    if data.get("anthropic_api_key") and not data["anthropic_api_key"].startswith("•"):
        current["anthropic_api_key"] = data["anthropic_api_key"]
    save_settings(current)
    return jsonify({"ok": True})


@app.route("/api/sent-log", methods=["GET"])
def get_sent_log():
    return jsonify({"ok": True, "log": load_sent_log()})


@app.route("/api/case-studies", methods=["GET"])
def list_case_studies():
    result = []
    for label, fname in CASE_STUDIES:
        path = get_case_study_path(fname)
        result.append({"label": label, "filename": fname, "exists": path is not None})
    return jsonify({"ok": True, "case_studies": result})


# ── Reference Library routes ──────────────────────────────────────────────────

@app.route("/api/library", methods=["GET"])
def get_library():
    """Return full 3-tier library."""
    meta = load_library_meta()

    # Tier 1 — Overall Reference
    overall   = list_tier_files(OVERALL_DIR,   meta.get("overall",   []))

    # Tier 2 — Follow-up 1 (sector decks — fixed slots + free uploads)
    fu1_fixed = []
    for sector_key, fname in SECTOR_DECKS.items():
        path = os.path.join(FOLLOWUP1_DIR, fname)
        info = file_info(path, fname)
        info["sector_key"]   = sector_key
        info["sector_label"] = SECTOR_LABELS.get(sector_key, sector_key)
        info["label"]        = SECTOR_LABELS.get(sector_key, fname)
        info["fixed"]        = True
        fu1_fixed.append(info)
    fu1_extra = [f for f in list_tier_files(FOLLOWUP1_DIR, meta.get("followup1", []))
                 if f["filename"] not in SECTOR_DECKS.values()]
    followup1 = fu1_fixed + fu1_extra

    # Tier 3 — Follow-up 2 (case studies — fixed slots + free uploads)
    fu2_fixed = []
    for label, fname in CASE_STUDIES:
        path = os.path.join(FOLLOWUP2_DIR, fname)
        info = file_info(path, fname)
        info["label"] = label
        info["fixed"] = True
        fu2_fixed.append(info)
    fu2_extra = [f for f in list_tier_files(FOLLOWUP2_DIR, meta.get("followup2", []))
                 if f["filename"] not in {fname for _, fname in CASE_STUDIES}]
    followup2 = fu2_fixed + fu2_extra

    return jsonify({"ok": True,
                    "overall":   overall,
                    "followup1": followup1,
                    "followup2": followup2})


def _save_tier_file(tier, f, label, desc):
    """Save an uploaded file to the correct tier directory and update meta."""
    dirs = {"overall": OVERALL_DIR, "followup1": FOLLOWUP1_DIR, "followup2": FOLLOWUP2_DIR}
    folder = dirs.get(tier)
    if not folder:
        return None, "Unknown tier"
    safe = re.sub(r"[^\w.\-]", "_", f.filename)
    dest = os.path.join(folder, safe)
    f.save(dest)
    meta = load_library_meta()
    tier_list = meta.setdefault(tier, [])
    tier_list[:] = [x for x in tier_list if x["filename"] != safe]
    tier_list.append({
        "filename":    safe,
        "label":       label or f.filename,
        "description": desc,
        "uploaded_at": datetime.now().strftime("%Y-%m-%d %H:%M"),
    })
    save_library_meta(meta)
    return file_info(dest, safe), None


@app.route("/api/library/<tier>/upload", methods=["POST"])
def upload_to_tier(tier):
    """Upload a file to any tier: overall | followup1 | followup2."""
    if tier not in ("overall", "followup1", "followup2"):
        return jsonify({"ok": False, "error": "Unknown tier"}), 400
    f     = request.files.get("file")
    label = request.form.get("label", "")
    desc  = request.form.get("description", "")
    if not f or not f.filename:
        return jsonify({"ok": False, "error": "No file provided"}), 400
    info, err = _save_tier_file(tier, f, label, desc)
    if err:
        return jsonify({"ok": False, "error": err}), 400
    return jsonify({"ok": True, "message": f"Uploaded to {tier}", **info})


@app.route("/api/library/<tier>/delete/<filename>", methods=["DELETE"])
def delete_from_tier(tier, filename):
    """Delete a file from a tier (only free-upload files, not fixed slots)."""
    dirs = {"overall": OVERALL_DIR, "followup1": FOLLOWUP1_DIR, "followup2": FOLLOWUP2_DIR}
    folder = dirs.get(tier)
    if not folder:
        return jsonify({"ok": False, "error": "Unknown tier"}), 400
    safe = re.sub(r"[^\w.\-]", "_", filename)
    path = os.path.join(folder, safe)
    if os.path.exists(path):
        os.remove(path)
    meta = load_library_meta()
    tier_list = meta.get(tier, [])
    meta[tier] = [x for x in tier_list if x["filename"] != safe]
    save_library_meta(meta)
    return jsonify({"ok": True})


@app.route("/api/library/followup1/deck/<sector_key>", methods=["POST"])
def replace_deck(sector_key):
    """Replace a specific sector deck (fixed slot in followup1)."""
    if sector_key not in SECTOR_DECKS:
        return jsonify({"ok": False, "error": "Unknown sector"}), 400
    f = request.files.get("file")
    if not f or not f.filename:
        return jsonify({"ok": False, "error": "No file"}), 400
    fname = SECTOR_DECKS[sector_key]
    dest  = os.path.join(FOLLOWUP1_DIR, fname)
    f.save(dest)
    return jsonify({"ok": True, "message": f"Deck updated: {SECTOR_LABELS[sector_key]}",
                    **file_info(dest, fname)})


@app.route("/api/library/followup2/case-study/<int:index>", methods=["POST"])
def replace_case_study(index):
    """Replace a specific case study (fixed slot in followup2)."""
    if index < 0 or index >= len(CASE_STUDIES):
        return jsonify({"ok": False, "error": "Invalid index"}), 400
    f = request.files.get("file")
    if not f or not f.filename:
        return jsonify({"ok": False, "error": "No file"}), 400
    _, fname = CASE_STUDIES[index]
    dest = os.path.join(FOLLOWUP2_DIR, fname)
    f.save(dest)
    return jsonify({"ok": True, "message": f"Case study {index+1} updated",
                    **file_info(dest, fname)})


@app.route("/api/library/download/<tier>/<filename>")
def download_library_file(tier, filename):
    """Download any library file."""
    dirs = {"overall": OVERALL_DIR, "followup1": FOLLOWUP1_DIR, "followup2": FOLLOWUP2_DIR}
    folder = dirs.get(tier)
    if not folder:
        return "Not found", 404
    safe = re.sub(r"[^\w.\-]", "_", filename)
    path = os.path.join(folder, safe)
    if not os.path.exists(path):
        return "File not found", 404
    return send_file(path, as_attachment=True)


@app.route("/api/health")
def health():
    return jsonify({"ok": True, "service": "FEAAM Priority Sender"})


if __name__ == "__main__":
    if not IS_CLOUD:
        _seed_local_files()   # copy local decks/case-studies into library on first run
    port  = int(os.getenv("PORT", 5055))
    debug = not IS_CLOUD
    print(f"\n  FEAAM Priority Sender  ->  http://localhost:{port}\n")
    app.run(debug=debug, port=port, host="0.0.0.0")
